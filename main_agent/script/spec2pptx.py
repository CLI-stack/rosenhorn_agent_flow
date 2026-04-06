from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN 
from pptx.dml.color import RGBColor
import argparse
import re
import random
from collections import Counter


def read_spec(specFile,pptx):
    spec = open(specFile,'r')
    # Create a presentation object
    presentation = Presentation()
    presentation.slide_height = Inches(9)
    presentation.slide_width = Inches(16)
    isText = 0
    isTitile = 0
    tableStart = 0
    isPic = 0
    text_list = []
    pic_list = []
    for line in spec:
        se = re.search('^#text#',line)
        if se:
            isText = 1
            continue
        
        if re.search("#pic#",line):
            isPic = 1

        se = re.search('^#table#',line)
        if se:
            table_list = []
            tableStart = 1
            continue

        se = re.search('^#table end#',line)
        if se:
            if tableStart == 1:
                tableStart = 0
                # Define table data
                table_width = 15.0
                # Set column widths according to number of words, 7 character need 1 inches
                inches_c = 7
                rows = len(table_list)
                cols = len(table_list[0].split('|'))
                # table offset
                left = Inches(0.5)
                top = Inches(2.0)
                width = Inches(table_width)
                height = Inches(0.8)
                # total column per slide
                cols_slide = 10 
                # Add a table to the slide
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                cols_orig_width = {}
                cols_orig_width_sum = 0
                cols_i_max = 0
                for line in table_list:
                    for i in range(len(line.split('|'))):
                        for sub_line in line.split('|')[i].split('\n'):
                            if re.search('/proj/|http|/home/',sub_line):
                                if  re.search('.log|.html|.png|http',sub_line):
                                    lable = sub_line.split('/')[-1]
                                    lable = lable.split('.')[0]
                                    sub_line = lable 

                            if i in cols_orig_width: 
                                if cols_orig_width[i] < round(len(sub_line)/inches_c,1):
                                    cols_orig_width[i] = round(len(sub_line)/inches_c,1)
                                if cols_orig_width[i] > max(cols_orig_width.values()):
                                    cols_i_max = i
                            else:
                                cols_orig_width[i] = round(len(sub_line)/inches_c,1)
                                cols_i_max = i
                        
                print(sum(cols_orig_width.values()))

                if sum(cols_orig_width.values()) < table_width:
                    for i in cols_orig_width:
                        table.columns[i].width = Inches(cols_orig_width[i])
                else:
                    for i in cols_orig_width:
                        if i == cols_i_max:
                            table.columns[i].width = Inches(table_width - (sum(cols_orig_width.values()) - cols_orig_width[i]))
                        else:
                            table.columns[i].width = Inches(cols_orig_width[i])

                rows = 0
                for line in table_list:
                    line = re.sub('\n','',line)
                    cols = 0
                    for cell in line.split('|'):
                        print("#cell",cell)
                        if len(cell.split(';')) > 1 :
                            print("use list")
                            cl = table.cell(rows, cols)
                            text_frame = cl.text_frame
                            for t in cell.split(';'):
                                if re.search('\S+',t):
                                    if re.search('/proj/|http|/home/',t):
                                        t = re.sub(' ','',t)
                                        if  re.search('.log|.html|.png|http',t):
                                            lable = t.split('/')[-1]
                                            lable = lable.split('.')[0]
                                            print(lable,t)
                                            if  re.search('http',t):
                                                p = text_frame.add_paragraph()
                                                r = p.add_run()
                                                r.text = lable
                                                r.hyperlink.address = t
                                                #data_ul += li(a(lable,href=t,align="left"))
                                            else:
                                                p = text_frame.add_paragraph()
                                                r = p.add_run()
                                                r.text = lable
                                                r.hyperlink.address = 'http://logviewer-atl.amd.com/'+t
                                                #data_ul += li(a(lable,href='http://logviewer-atl.amd.com/'+t,align="left"))
                                        else:
                                            p = text_frame.add_paragraph()
                                            r = p.add_run()
                                            r.text = lable
                                            r.hyperlink.address = 'http://logviewer-atl.amd.com/'+t

                                            #data_ul += li(a(t,href='http://logviewer-atl.amd.com/'+t,align="left"))
                                    else:
                                        p = text_frame.add_paragraph()
                                        r = p.add_run()
                                        r.text = re.sub("^ ","",t)

                            for paragraph in text_frame.paragraphs:
                                if paragraph.text == '':
                                    text_frame._element.remove(paragraph._element)

                                        #data_ul += li(t,align="left")
                            pass
                            #data_tr+=td(data_ul,align="left")
                        else:
                            #data_td = td()

                            if len(cell.split('::')) == 2:
                                color = cell.split('::')[1]
                                # transfer dex color(#FF0000) to dec rgb(255,0,0)
                                color = re.sub("#","",color)
                                print(color,color[0:2],color[2:4])
                                r = int(color[0:2],16)
                                g = int(color[2:4],16)
                                b = int(color[4:6],16)
                                
                                cell_text = cell.split('::')[0]
                            else:
                                cell_text = cell.split('::')[0]
                                color = ""

                            if re.search('/proj/|http|/home/',cell_text):
                                if  re.search('.log|.html|.png|http',cell_text):
                                    lable = cell_text.split('/')[-1]
                                    lable = lable.split('.')[0]
                                    if  re.search('http',cell_text):
                                        cl = table.cell(rows, cols)
                                        text_frame = cl.text_frame
                                        p = text_frame.add_paragraph()
                                        r = p.add_run()
                                        r.text = lable
                                        print("# lable",lable)
                                        cell_text = re.sub(" ","",cell_text)
                                        r.hyperlink.address = cell_text
                                        for paragraph in text_frame.paragraphs:
                                            if paragraph.text == '':
                                                text_frame._element.remove(paragraph._element)
                
                                        print(cell_text)
                                    else:
                                        #data_tr+=td(a(lable,href='http://logviewer-atl.amd.com/'+cell_text),style=style)
                                        cl = table.cell(rows, cols)
                                        text_frame = cl.text_frame
                                        p = text_frame.add_paragraph()
                                        r = p.add_run()
                                        r.text = lable
                                        print("# lable:",lable)
                                        cell_text = re.sub(" ","",cell_text)
                                        r.hyperlink.address = 'http://logviewer-atl.amd.com/'+cell_text
                                        for paragraph in text_frame.paragraphs:
                                            if paragraph.text == '':
                                                text_frame._element.remove(paragraph._element)

                                        print('http://logviewer-atl.amd.com/'+cell_text)
                                else:
                                    table.cell(rows, cols).text = cell_text
                                    #data_tr+=td(a(cell_text,href='http://logviewer-atl.amd.com/'+cell_text),style=style)
                            else:
                                table.cell(rows, cols).text = cell_text
                                if re.search("\S",color):
                                    table.cell(rows, cols).fill.solid()
                                    print("# fill color",r,g,b)
                                    table.cell(rows, cols).fill.fore_color.rgb = RGBColor(r,g,b)

                        cols = cols + 1
                    rows = rows + 1
            continue

        se = re.search('^#title#',line)
        if se:
            isTitile = 1
            continue

        if tableStart == 1:
            table_list.append(line)

        if isTitile == 1:  
            if re.search("summary",line,re.I):
                # Create a slide layout
                bullet_slide_layout = presentation.slide_layouts[1]
                # Add a slide with a title and content layout
                slide = presentation.slides.add_slide(bullet_slide_layout)
                # Add a title
                print("# add title",line)
                title = slide.shapes.title
                title.text = line
                title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT 
            else:
                # Add another slide for the table
                print("# Add another slide for the table")
                slide_layout = presentation.slide_layouts[5]
                slide = presentation.slides.add_slide(slide_layout)

                # Add a title
                title = slide.shapes.title
                title.text = line
                title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                print(slide)
                #add_table(slide)

            isTitile = 0
            continue

        if isText == 1:
            if re.search('\S',line):
                line = re.sub('\n','',line)
                text_list.append(line)
            else:
                # Add content
                print("# add text")
                content = slide.placeholders[1]
                content.text = "\n".join(text_list)
                text_list = []
                isText = 0

        if isPic == 1:
            if re.search('\S',line):
                line = re.sub('\n','',line)
                pic_list.append(line)
            else:
                print("# add picture")
                for pic in pic_list:
                    left = Inches(0.5)
                    top = Inches(0.5)
                    width = Inches(5.0)
                    height = Inches(5.0)
                    slide.shapes.add_picture(pic,left,top,width,height)
                isPic = 0
            
    

    print(slide)
    presentation.save(pptx)

def add_table(slide):
    # Define table data
    rows = 3
    cols = 3
    left = Inches(2.0)
    top = Inches(2.0)
    width = Inches(6.0)
    height = Inches(0.8)

    # Add a table to the slide
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(1.0)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(3.0)

    # Write column headings
    table.cell(0, 0).text = "Tile"
    table.cell(0, 1).text = "Issue"
    table.cell(0, 2).text = "Comment"

    # Write body cells
    table.cell(1, 0).text = "df_tcdxa_t"
    table.cell(1, 1).text = "ports are misaligned"
    table.cell(1, 2).text = "fcfp will fix"

    table.cell(2, 0).text = "umc_umcch_t"
    table.cell(2, 1).text = "congestion is high"
    table.cell(2, 2).text = "fcfp will increase the height"


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='spec to html')
    parser.add_argument('--spec',type=str, default = "None",required=True,help="spec file")
    parser.add_argument('--pptx',type=str, default = "None",required=True,help="pptx file")
    args = parser.parse_args()

    read_spec(args.spec,args.pptx)

